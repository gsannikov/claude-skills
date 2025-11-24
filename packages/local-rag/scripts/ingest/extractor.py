from pathlib import Path
from pypdf import PdfReader
from pdf2image import convert_from_path
from PIL import Image
import os

from .ocr import run_ocr

USE_OCR = os.getenv("OCR_ENABLED","true").lower()=="true"
OCR_MAX_PAGES = int(os.getenv("OCR_MAX_PAGES","120"))
OCR_PAGE_DPI = int(os.getenv("OCR_PAGE_DPI","200"))
TEXT_EXTS = {".txt",".md",".json",".yml",".yaml",".py",".js",".ts",".c",".cpp",".h",".sh"}
IMAGE_EXTS = {".png",".jpg",".jpeg",".tiff",".webp"}

def _read_docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    return "\n".join([para.text for para in d.paragraphs])

def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    return "\n".join(shp.text for s in prs.slides for shp in s.shapes if hasattr(shp,"text"))

def _read_xlsx(p: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    return "\n".join("\t".join("" if v is None else str(v) for v in row) for ws in wb for row in ws.iter_rows(values_only=True))

def read_text_with_ocr(p: Path) -> str:
    ext = p.suffix.lower()
    if ext in TEXT_EXTS:
        return p.read_text(errors="ignore")

    if ext == ".docx":
        return _read_docx(p)

    if ext == ".pptx":
        return _read_pptx(p)

    if ext == ".xlsx":
        return _read_xlsx(p)

    if ext == ".pdf":
        try:
            r = PdfReader(str(p))
            texts = [(pg.extract_text() or "") for pg in r.pages]
            joined = "\n".join(texts).strip()
            if not USE_OCR or (len(texts) and (sum(len(t) for t in texts)/max(1,len(texts)) >= 100)):
                return joined
            pages = min(len(r.pages), OCR_MAX_PAGES)
        except Exception:
            pages = None

        if USE_OCR:
            imgs = convert_from_path(str(p), dpi=OCR_PAGE_DPI, first_page=1, last_page=(pages or 1))
            return run_ocr(imgs)
        return joined

    if ext in IMAGE_EXTS and USE_OCR:
        img = Image.open(str(p)).convert("RGB")
        return run_ocr([img])

    return ""
