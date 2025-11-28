from pathlib import Path
from pypdf import PdfReader
from pdf2image import convert_from_path
from PIL import Image

from ..settings import LocalRagSettings, get_settings
from .ocr import run_ocr

# Text-based files (read directly)
TEXT_EXTS = {
    ".txt", ".md", ".mdx", ".mdc", ".log",
    ".json", ".jsonc", ".yml", ".yaml",
    ".py", ".js", ".ts", ".tsx", ".jsx",
    ".c", ".cpp", ".h", ".hpp", ".sh",
    ".html", ".xhtml", ".css", ".scss",
    ".swift", ".go", ".rs", ".java",
}

# Image files (require OCR)
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".webp"}

def _read_docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    return "\n".join([para.text for para in d.paragraphs])

def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    return "\n".join(shp.text for s in prs.slides for shp in s.shapes if hasattr(shp,"text"))

def _read_xlsx(p: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    return "\n".join("\t".join("" if v is None else str(v) for v in row) for ws in wb for row in ws.iter_rows(values_only=True))

def _read_doc(p: Path) -> str:
    """Read legacy .doc files using antiword or textract fallback."""
    import subprocess
    try:
        # Try antiword first (brew install antiword)
        result = subprocess.run(["antiword", str(p)], capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Fallback: try textract if available
    try:
        import textract
        return textract.process(str(p)).decode('utf-8', errors='ignore')
    except ImportError:
        pass

    return ""

def read_text_with_ocr(p: Path, settings: LocalRagSettings | None = None) -> str:
    settings = settings or get_settings()
    ext = p.suffix.lower()
    if ext in TEXT_EXTS:
        return p.read_text(errors="ignore")

    if ext == ".docx":
        return _read_docx(p)

    if ext == ".doc":
        return _read_doc(p)

    if ext == ".pptx":
        return _read_pptx(p)

    if ext == ".xlsx":
        return _read_xlsx(p)

    if ext == ".pdf":
        try:
            r = PdfReader(str(p))
            texts = [(pg.extract_text() or "") for pg in r.pages]
            joined = "\n".join(texts).strip()
            if not settings.ocr_enabled or (len(texts) and (sum(len(t) for t in texts)/max(1,len(texts)) >= 100)):
                return joined
            pages = min(len(r.pages), settings.ocr_max_pages)
        except Exception as e:
            print(f"PDF text extraction failed for {p.name}: {e}")
            pages = None
        if settings.ocr_enabled:
            imgs = convert_from_path(str(p), dpi=settings.ocr_page_dpi, first_page=1, last_page=(pages or 1))
            return run_ocr(imgs, settings=settings)
        return joined

    if ext in IMAGE_EXTS and settings.ocr_enabled:
        img = Image.open(str(p)).convert("RGB")
        return run_ocr([img], settings=settings)

    return ""
