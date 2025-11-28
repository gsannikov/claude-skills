import logging
from pathlib import Path
from pypdf import PdfReader
from pdf2image import convert_from_path
from PIL import Image

# Raise the PIL pixel limit to avoid DecompressionBomb warnings on moderate images.
Image.MAX_IMAGE_PIXELS = 200_000_000

from ..settings import LocalRagSettings, get_settings
from .ocr import run_ocr, _resolve_tesseract_lang

# Text-based files (read directly)
TEXT_EXTS = {
    ".txt", ".md", ".mdx", ".mdc", ".log",
    ".json", ".jsonc", ".yml", ".yaml",
    ".py", ".js", ".ts", ".tsx", ".jsx",
    ".c", ".cpp", ".h", ".hpp", ".sh",
    ".html", ".xhtml", ".css", ".scss",
    ".swift", ".go", ".rs", ".java",
}

# Image files (require OCR)
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".webp"}

def _configure_pdf_logging():
    """Keep noisy pypdf warnings out of stdout during ingestion."""
    logging.getLogger("pypdf").setLevel(logging.ERROR)


def _read_docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    return "\n".join([para.text for para in d.paragraphs])

def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    return "\n".join(shp.text for s in prs.slides for shp in s.shapes if hasattr(shp,"text"))

def _read_xlsx(p: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    return "\n".join("\t".join("" if v is None else str(v) for v in row) for ws in wb for row in ws.iter_rows(values_only=True))

def _read_doc(p: Path) -> str:
    """Read legacy .doc files using antiword or textract fallback."""
    import subprocess
    try:
        # Try antiword first (brew install antiword)
        result = subprocess.run(["antiword", str(p)], capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Fallback: try textract if available
    try:
        import textract
        return textract.process(str(p)).decode('utf-8', errors='ignore')
    except ImportError:
        pass

    return ""

def read_text_with_ocr(p: Path, settings: LocalRagSettings | None = None) -> str:
    settings = settings or get_settings()
    ext = p.suffix.lower()
    if ext in TEXT_EXTS:
        return p.read_text(errors="ignore")

    if ext == ".docx":
        return _read_docx(p)

    if ext == ".doc":
        return _read_doc(p)

    if ext == ".pptx":
        return _read_pptx(p)

    if ext == ".xlsx":
        return _read_xlsx(p)

    if ext == ".pdf":
        print(f"PDF: Reading {p.name}...", flush=True)
        try:
            _configure_pdf_logging()
            r = PdfReader(str(p), strict=False)
            print(f"PDF: Found {len(r.pages)} pages", flush=True)
            texts = [(pg.extract_text() or "") for pg in r.pages]
            joined = "\n".join(texts).strip()
            avg_text_len = sum(len(t) for t in texts)/max(1,len(texts))
            print(f"PDF: Average {avg_text_len:.0f} chars/page", flush=True)
            
            if not settings.ocr_enabled or (len(texts) and avg_text_len >= 100):
                print(f"PDF: Using extracted text (OCR={settings.ocr_enabled}, avg_len={avg_text_len:.0f})", flush=True)
                return joined
            pages = min(len(r.pages), settings.ocr_max_pages)
            print(f"PDF: Will OCR first {pages} pages", flush=True)
        except Exception as e:
            print(f"PDF text extraction failed for {p.name}: {e}", flush=True)
            pages = None
        if settings.ocr_enabled:
            print(f"PDF: OCR required for {p.name}...", flush=True)
            try:
                import ocrmypdf
                import tempfile
                
                # Map language(s) to Tesseract codes
                ocr_lang = _resolve_tesseract_lang(settings.ocr_lang)
                
                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tf:
                    print(f"PDF: Running OCRmyPDF (lang={ocr_lang})...", flush=True)
                    try:
                        ocrmypdf.ocr(
                            str(p), 
                            tf.name, 
                            language=ocr_lang, 
                            force_ocr=True, 
                            progress_bar=False,
                            optimize=1,
                            skip_text=False
                        )
                    except ocrmypdf.exceptions.MissingDependencyError as e:
                        print(f"PDF OCR missing dependency: {e}", flush=True)
                        print("Install system deps: tesseract-ocr, qpdf, ghostscript, poppler-utils", flush=True)
                        raise
                    
                    # Read text from the OCR'd PDF
                    print(f"PDF: OCR complete, extracting text...", flush=True)
                    r_ocr = PdfReader(tf.name)
                    ocr_text = "\n".join([(pg.extract_text() or "") for pg in r_ocr.pages])
                    print(f"PDF: Extracted {len(ocr_text)} chars from OCR", flush=True)
                    return ocr_text.strip()
            except ImportError:
                print("Error: ocrmypdf not installed. Please install it with 'pip install ocrmypdf'", flush=True)
            except Exception as e:
                print(f"PDF OCR failed: {e}", flush=True)
                # Fallback to image-based OCR if ocrmypdf fails (e.g. missing tesseract)
                print("PDF: Falling back to image-based OCR...", flush=True)
                imgs = convert_from_path(str(p), dpi=settings.ocr_page_dpi, first_page=1, last_page=(pages or 1))
                return run_ocr(imgs, settings=settings)
        return joined

    if ext in IMAGE_EXTS and settings.ocr_enabled:
        img = Image.open(str(p))
        # Skip extremely large images that could trigger PIL DecompressionBomb
        if img.width * img.height > 80_000_000:
            print(f"Skipping oversized image (>{80_000_000} px): {p.name}")
            return ""
        img = img.convert("RGB")
        return run_ocr([img], settings=settings)

    return ""
